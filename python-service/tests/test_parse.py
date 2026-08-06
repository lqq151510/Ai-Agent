import pytest
from fastapi.testclient import TestClient
from app.main import app
import tempfile
import os

client = TestClient(app)

def test_parse_document_success(monkeypatch):
    from app.api.routers import parse

    def mock_parse_document(file_path):
        return f"# Converted {os.path.basename(file_path)}", {"sourceFormat": "txt"}

    monkeypatch.setattr(parse, "parse_document", mock_parse_document)

    # Create a dummy file
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp.write(b"Hello World")
        tmp_name = tmp.name

    try:
        with open(tmp_name, "rb") as f:
            response = client.post("/parse", files={"file": ("test.txt", f, "text/plain")})
        
        assert response.status_code == 200
        data = response.json()
        assert data["filename"] == "test.txt"
        assert "# Converted" in data["markdown"]
    finally:
        if os.path.exists(tmp_name):
            os.remove(tmp_name)

def test_parse_document_missing_file():
    response = client.post("/parse")
    assert response.status_code == 422 # Unprocessable Entity (Missing file field)


@pytest.mark.parametrize(
    ("filename", "media_type", "source_format"),
    [
        (
            "learning-plan.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
        ),
        (
            "architecture.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
        ),
    ],
)
def test_parse_document_accepts_modern_office_extensions(
    monkeypatch, filename, media_type, source_format
):
    from app.api.routers import parse

    monkeypatch.setattr(
        parse,
        "parse_document",
        lambda file_path: ("# Converted Office document", {"sourceFormat": source_format}),
    )

    response = client.post(
        "/parse",
        files={"file": (filename, b"office bytes", media_type)},
    )

    assert response.status_code == 200
    assert response.json()["source_format"] == source_format


def test_parse_document_converts_real_modern_office_files():
    from docx import Document
    from pptx import Presentation
    from app.services.parser_service import parse_document

    with tempfile.TemporaryDirectory() as directory:
        docx_path = os.path.join(directory, "learning-plan.docx")
        document = Document()
        document.add_heading("RAG Learning Plan", level=1)
        document.add_paragraph("Practice retrieval and reranking every week.")
        document.core_properties.title = "RAG Learning Plan"
        document.save(docx_path)

        docx_markdown, docx_metadata = parse_document(docx_path)

        assert "RAG Learning Plan" in docx_markdown
        assert docx_metadata["sourceFormat"] == "docx"
        assert docx_metadata["title"] == "RAG Learning Plan"
        assert int(docx_metadata["paragraphCount"]) >= 2

        pptx_path = os.path.join(directory, "architecture.pptx")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Knowledge Desk Architecture"
        slide.placeholders[1].text = "Import, organize, search, and review."
        presentation.core_properties.title = "Knowledge Desk Architecture"
        presentation.save(pptx_path)

        pptx_markdown, pptx_metadata = parse_document(pptx_path)

        assert "Knowledge Desk Architecture" in pptx_markdown
        assert pptx_metadata["sourceFormat"] == "pptx"
        assert pptx_metadata["title"] == "Knowledge Desk Architecture"
        assert pptx_metadata["slideCount"] == "1"
